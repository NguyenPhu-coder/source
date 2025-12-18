"""
Content Ingestion Agent - Multi-modal content processing with ML models
Handles PDF/DOCX/PPTX ingestion, analysis, and knowledge graph integration
"""

import asyncio
import hashlib
import io
import json
import os
import re
import shutil
import tempfile
import uuid
from contextlib import asynccontextmanager
from datetime import datetime
from enum import Enum
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import cv2
import httpx
import numpy as np
import pytesseract
import structlog
import torch
import yaml
from celery import Celery
from docx import Document
from fastapi import FastAPI, File, HTTPException, UploadFile, WebSocket, status
from fastapi.responses import JSONResponse
from pdf2image import convert_from_path
from PIL import Image
from prometheus_client import Counter, Histogram, generate_latest
from pydantic import BaseModel, Field, validator
from PyPDF2 import PdfReader
from pptx import Presentation
from tenacity import retry, stop_after_attempt, wait_exponential
from transformers import (
    AutoModelForSequenceClassification,
    AutoTokenizer,
    AutoModelForTokenClassification,
    pipeline,
)

# Configure structured logging
structlog.configure(
    processors=[
        structlog.contextvars.merge_contextvars,
        structlog.processors.add_log_level,
        structlog.processors.TimeStamper(fmt="iso"),
        structlog.processors.JSONRenderer(),
    ]
)

logger = structlog.get_logger()

# Prometheus metrics
INGESTION_COUNT = Counter(
    "ingestion_total", "Total ingestions", ["status", "file_type"]
)
INGESTION_DURATION = Histogram(
    "ingestion_duration_seconds", "Ingestion duration", ["file_type"]
)
EXTRACTION_COUNT = Counter(
    "extraction_total", "Content extractions", ["content_type"]
)
MODEL_INFERENCE = Histogram(
    "model_inference_seconds", "Model inference time", ["model_name"]
)
GRAPH_API_CALLS = Counter(
    "graph_api_calls_total", "Knowledge graph API calls", ["operation", "status"]
)


class FileFormat(str, Enum):
    """Supported file formats"""
    PDF = "pdf"
    DOCX = "docx"
    PPTX = "pptx"


class ContentType(str, Enum):
    """Content types"""
    TEXT = "text"
    IMAGE = "image"
    TABLE = "table"
    EQUATION = "equation"
    DIAGRAM = "diagram"


class ProcessingStatus(str, Enum):
    """Processing status"""
    PENDING = "pending"
    PROCESSING = "processing"
    COMPLETED = "completed"
    FAILED = "failed"


class IngestionRequest(BaseModel):
    """Request model for file ingestion"""
    file_name: str
    file_type: FileFormat
    metadata: Dict[str, Any] = Field(default_factory=dict)
    
    @validator("metadata")
    def validate_metadata(cls, v):
        """Validate metadata structure"""
        allowed_keys = {"subject", "grade", "author", "language", "tags"}
        if not all(k in allowed_keys for k in v.keys()):
            raise ValueError(f"Invalid metadata keys. Allowed: {allowed_keys}")
        return v


class IngestionResponse(BaseModel):
    """Response model for ingestion"""
    job_id: str
    status: ProcessingStatus
    file_id: str
    message: str
    created_at: datetime


class ExtractionResponse(BaseModel):
    """Response model for content extraction"""
    file_id: str
    content: Dict[str, Any]
    concepts: List[Dict[str, Any]]
    relationships: List[Dict[str, Any]]
    classification: Dict[str, Any]
    extracted_at: datetime


class StatusResponse(BaseModel):
    """Response model for status check"""
    job_id: str
    status: ProcessingStatus
    progress: float
    message: str
    result: Optional[Dict[str, Any]] = None


class ModelManager:
    """Manages ML model loading and inference"""
    
    def __init__(self, config: Dict[str, Any]):
        self.config = config
        self.models = {}
        self.tokenizers = {}
        self.pipelines = {}
        self.logger = structlog.get_logger()
        
    def load_models(self):
        """Load all required ML models"""
        self.logger.info("loading_models_started")
        
        try:
            # Load DistilBERT for classification
            distilbert_config = self.config["models"]["distilbert"]
            model_path = distilbert_config["model_path"]
            num_labels = distilbert_config["num_labels"]
            
            self.logger.info("loading_distilbert", path=model_path)
            self.tokenizers["distilbert"] = AutoTokenizer.from_pretrained(model_path)
            self.models["distilbert"] = AutoModelForSequenceClassification.from_pretrained(
                model_path,
                num_labels=num_labels
            )
            
            # Load NER model for concept extraction
            self.logger.info("loading_ner_model")
            self.pipelines["ner"] = pipeline(
                "ner",
                model="dslim/bert-base-NER",
                aggregation_strategy="simple"
            )
            
            # Load zero-shot classification for relationships
            self.logger.info("loading_zero_shot_classifier")
            self.pipelines["zero_shot"] = pipeline(
                "zero-shot-classification",
                model="facebook/bart-large-mnli"
            )
            
            self.logger.info("models_loaded_successfully")
            
        except Exception as e:
            self.logger.error("model_loading_failed", error=str(e))
            raise
    
    def classify_content(self, text: str) -> Dict[str, Any]:
        """
        Classify content using DistilBERT
        
        Args:
            text: Input text
            
        Returns:
            Classification results with subject, difficulty, bloom's level
        """
        start_time = datetime.utcnow()
        
        try:
            # Tokenize and classify
            inputs = self.tokenizers["distilbert"](
                text,
                padding=True,
                truncation=True,
                max_length=512,
                return_tensors="pt"
            )
            
            with torch.no_grad():
                outputs = self.models["distilbert"](**inputs)
                predictions = torch.nn.functional.softmax(outputs.logits, dim=-1)
            
            # Get top predictions
            top_k = 3
            probs, indices = torch.topk(predictions[0], top_k)
            
            subjects = self.config["models"]["distilbert"]["subjects"]
            
            results = {
                "primary_subject": subjects[indices[0].item()],
                "confidence": float(probs[0].item()),
                "top_predictions": [
                    {
                        "subject": subjects[idx.item()],
                        "confidence": float(prob.item())
                    }
                    for prob, idx in zip(probs, indices)
                ],
                "difficulty": self._estimate_difficulty(text),
                "blooms_level": self._estimate_blooms_level(text)
            }
            
            duration = (datetime.utcnow() - start_time).total_seconds()
            MODEL_INFERENCE.labels(model_name="distilbert").observe(duration)
            
            return results
            
        except Exception as e:
            self.logger.error("classification_failed", error=str(e))
            raise
    
    def extract_concepts(self, text: str, method: str = "ner") -> List[Dict[str, Any]]:
        """
        Extract concepts using NER
        
        Args:
            text: Input text
            method: Extraction method (ner, keyword, hybrid)
            
        Returns:
            List of extracted concepts with metadata
        """
        start_time = datetime.utcnow()
        
        try:
            concepts = []
            
            if method in ["ner", "hybrid"]:
                # NER-based extraction
                entities = self.pipelines["ner"](text)
                
                for entity in entities:
                    if entity["score"] >= self.config["extraction"]["min_confidence"]:
                        concepts.append({
                            "text": entity["word"],
                            "type": entity["entity_group"],
                            "confidence": entity["score"],
                            "start": entity["start"],
                            "end": entity["end"],
                            "method": "ner"
                        })
            
            if method in ["keyword", "hybrid"]:
                # Keyword-based extraction
                keywords = self._extract_keywords(text)
                for keyword in keywords:
                    concepts.append({
                        "text": keyword["text"],
                        "type": "KEYWORD",
                        "confidence": keyword["score"],
                        "method": "keyword"
                    })
            
            # Deduplicate
            seen = set()
            unique_concepts = []
            for concept in concepts:
                key = concept["text"].lower()
                if key not in seen:
                    seen.add(key)
                    unique_concepts.append(concept)
            
            duration = (datetime.utcnow() - start_time).total_seconds()
            MODEL_INFERENCE.labels(model_name="ner").observe(duration)
            
            self.logger.info(
                "concepts_extracted",
                count=len(unique_concepts),
                method=method
            )
            
            return unique_concepts
            
        except Exception as e:
            self.logger.error("concept_extraction_failed", error=str(e))
            raise
    
    def detect_relationships(
        self,
        concepts: List[Dict[str, Any]],
        text: str
    ) -> List[Dict[str, Any]]:
        """
        Detect relationships between concepts
        
        Args:
            concepts: List of extracted concepts
            text: Original text
            
        Returns:
            List of relationships with types and confidence
        """
        start_time = datetime.utcnow()
        
        try:
            relationships = []
            
            # Relationship templates
            rel_templates = [
                "prerequisite_of",
                "related_to",
                "example_of",
                "part_of",
                "causes",
                "defines"
            ]
            
            # Check concept pairs
            for i, concept1 in enumerate(concepts):
                for concept2 in concepts[i+1:]:
                    # Check if concepts appear in same sentence
                    if self._in_same_context(concept1, concept2, text):
                        # Classify relationship type
                        context = self._get_context(concept1, concept2, text)
                        
                        result = self.pipelines["zero_shot"](
                            context,
                            candidate_labels=rel_templates
                        )
                        
                        if result["scores"][0] >= self.config["extraction"]["min_confidence"]:
                            relationships.append({
                                "from_concept": concept1["text"],
                                "to_concept": concept2["text"],
                                "relationship_type": result["labels"][0],
                                "confidence": result["scores"][0],
                                "context": context
                            })
            
            duration = (datetime.utcnow() - start_time).total_seconds()
            MODEL_INFERENCE.labels(model_name="zero_shot").observe(duration)
            
            self.logger.info("relationships_detected", count=len(relationships))
            
            return relationships
            
        except Exception as e:
            self.logger.error("relationship_detection_failed", error=str(e))
            raise
    
    def _estimate_difficulty(self, text: str) -> str:
        """Estimate content difficulty"""
        # Simple heuristic based on text complexity
        words = text.split()
        avg_word_len = sum(len(w) for w in words) / max(len(words), 1)
        sentences = text.split('.')
        avg_sentence_len = len(words) / max(len(sentences), 1)
        
        if avg_word_len < 5 and avg_sentence_len < 15:
            return "beginner"
        elif avg_word_len < 7 and avg_sentence_len < 25:
            return "intermediate"
        else:
            return "advanced"
    
    def _estimate_blooms_level(self, text: str) -> str:
        """Estimate Bloom's taxonomy level"""
        # Keyword-based heuristic
        keywords = {
            "remember": ["define", "list", "identify", "recall"],
            "understand": ["explain", "describe", "summarize", "interpret"],
            "apply": ["solve", "demonstrate", "use", "implement"],
            "analyze": ["compare", "contrast", "examine", "analyze"],
            "evaluate": ["assess", "critique", "judge", "evaluate"],
            "create": ["design", "construct", "create", "develop"]
        }
        
        text_lower = text.lower()
        for level, words in keywords.items():
            if any(word in text_lower for word in words):
                return level
        
        return "understand"  # Default
    
    def _extract_keywords(self, text: str, top_k: int = 10) -> List[Dict[str, Any]]:
        """Extract keywords using TF-IDF-like scoring"""
        # Simple keyword extraction
        words = re.findall(r'\b[a-z]{4,}\b', text.lower())
        word_freq = {}
        for word in words:
            word_freq[word] = word_freq.get(word, 0) + 1
        
        # Sort by frequency
        sorted_words = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)
        
        return [
            {"text": word, "score": freq / len(words)}
            for word, freq in sorted_words[:top_k]
        ]
    
    def _in_same_context(
        self,
        concept1: Dict[str, Any],
        concept2: Dict[str, Any],
        text: str,
        window: int = 100
    ) -> bool:
        """Check if concepts appear in same context window"""
        if "start" in concept1 and "start" in concept2:
            return abs(concept1["start"] - concept2["start"]) <= window
        return False
    
    def _get_context(
        self,
        concept1: Dict[str, Any],
        concept2: Dict[str, Any],
        text: str,
        window: int = 50
    ) -> str:
        """Get text context around concepts"""
        if "start" in concept1 and "start" in concept2:
            start = min(concept1["start"], concept2["start"]) - window
            end = max(concept1.get("end", 0), concept2.get("end", 0)) + window
            start = max(0, start)
            end = min(len(text), end)
            return text[start:end]
        return f"{concept1['text']} and {concept2['text']}"


class FileProcessor:
    """Handles file parsing and content extraction"""
    
    def __init__(self, config: Dict[str, Any]):
        self.config = config
        self.temp_dir = Path(config["file_processing"]["temp_dir"])
        self.output_dir = Path(config["file_processing"]["output_dir"])
        self.temp_dir.mkdir(parents=True, exist_ok=True)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.logger = structlog.get_logger()
    
    def validate_file(self, file_path: Path) -> bool:
        """
        Validate file format and size
        
        Args:
            file_path: Path to file
            
        Returns:
            True if valid
        """
        try:
            # Check file exists
            if not file_path.exists():
                raise ValueError(f"File not found: {file_path}")
            
            # Check file size
            max_size = self.config["file_processing"]["max_file_size_mb"] * 1024 * 1024
            if file_path.stat().st_size > max_size:
                raise ValueError(f"File too large: {file_path.stat().st_size} bytes")
            
            # Check format
            suffix = file_path.suffix.lower().lstrip('.')
            if suffix not in self.config["file_processing"]["supported_formats"]:
                raise ValueError(f"Unsupported format: {suffix}")
            
            self.logger.info("file_validated", path=str(file_path))
            return True
            
        except Exception as e:
            self.logger.error("file_validation_failed", error=str(e))
            raise
    
    def extract_content(self, file_path: Path) -> Dict[str, Any]:
        """
        Extract content from file based on format
        
        Args:
            file_path: Path to file
            
        Returns:
            Extracted content with text, images, tables
        """
        start_time = datetime.utcnow()
        
        try:
            suffix = file_path.suffix.lower().lstrip('.')
            
            if suffix == "pdf":
                content = self._extract_pdf(file_path)
            elif suffix == "docx":
                content = self._extract_docx(file_path)
            elif suffix == "pptx":
                content = self._extract_pptx(file_path)
            else:
                raise ValueError(f"Unsupported format: {suffix}")
            
            duration = (datetime.utcnow() - start_time).total_seconds()
            INGESTION_DURATION.labels(file_type=suffix).observe(duration)
            EXTRACTION_COUNT.labels(content_type="all").inc()
            
            self.logger.info(
                "content_extracted",
                file_type=suffix,
                pages=content.get("page_count", 0),
                duration=duration
            )
            
            return content
            
        except Exception as e:
            self.logger.error("content_extraction_failed", error=str(e))
            raise
    
    def _extract_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from PDF"""
        content = {
            "type": "pdf",
            "pages": [],
            "images": [],
            "tables": [],
            "equations": [],
            "metadata": {}
        }
        
        try:
            # Extract text
            reader = PdfReader(str(file_path))
            content["page_count"] = len(reader.pages)
            content["metadata"] = dict(reader.metadata) if reader.metadata else {}
            
            for page_num, page in enumerate(reader.pages):
                page_text = page.extract_text()
                
                # Extract equations (LaTeX patterns)
                equations = self._extract_equations(page_text)
                
                content["pages"].append({
                    "page_number": page_num + 1,
                    "text": page_text,
                    "word_count": len(page_text.split())
                })
                
                content["equations"].extend(equations)
            
            # Extract images
            images = convert_from_path(str(file_path))
            for idx, img in enumerate(images):
                img_path = self.temp_dir / f"{file_path.stem}_page_{idx}.png"
                img.save(img_path)
                
                # OCR on image
                ocr_text = pytesseract.image_to_string(img)
                
                content["images"].append({
                    "page": idx + 1,
                    "path": str(img_path),
                    "ocr_text": ocr_text,
                    "size": img.size
                })
            
            EXTRACTION_COUNT.labels(content_type="text").inc()
            EXTRACTION_COUNT.labels(content_type="image").inc(len(images))
            
        except Exception as e:
            self.logger.error("pdf_extraction_failed", error=str(e))
            raise
        
        return content
    
    def _extract_docx(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from DOCX"""
        content = {
            "type": "docx",
            "pages": [],
            "images": [],
            "tables": [],
            "metadata": {}
        }
        
        try:
            doc = Document(str(file_path))
            
            # Extract metadata
            content["metadata"] = {
                "author": doc.core_properties.author,
                "title": doc.core_properties.title,
                "created": str(doc.core_properties.created)
            }
            
            # Extract paragraphs
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)
            
            content["pages"].append({
                "page_number": 1,
                "text": "\n".join(full_text),
                "word_count": len(" ".join(full_text).split())
            })
            
            # Extract tables
            for table_idx, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                
                content["tables"].append({
                    "table_number": table_idx + 1,
                    "data": table_data,
                    "rows": len(table_data),
                    "cols": len(table_data[0]) if table_data else 0
                })
            
            content["page_count"] = 1
            
            EXTRACTION_COUNT.labels(content_type="text").inc()
            EXTRACTION_COUNT.labels(content_type="table").inc(len(doc.tables))
            
        except Exception as e:
            self.logger.error("docx_extraction_failed", error=str(e))
            raise
        
        return content
    
    def _extract_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from PPTX"""
        content = {
            "type": "pptx",
            "pages": [],
            "images": [],
            "tables": [],
            "metadata": {}
        }
        
        try:
            prs = Presentation(str(file_path))
            
            # Extract metadata
            content["metadata"] = {
                "author": prs.core_properties.author,
                "title": prs.core_properties.title,
                "slide_count": len(prs.slides)
            }
            
            # Extract slides
            for slide_idx, slide in enumerate(prs.slides):
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                
                content["pages"].append({
                    "page_number": slide_idx + 1,
                    "text": "\n".join(slide_text),
                    "word_count": len(" ".join(slide_text).split())
                })
            
            content["page_count"] = len(prs.slides)
            
            EXTRACTION_COUNT.labels(content_type="text").inc()
            
        except Exception as e:
            self.logger.error("pptx_extraction_failed", error=str(e))
            raise
        
        return content
    
    def _extract_equations(self, text: str) -> List[Dict[str, Any]]:
        """Extract LaTeX equations from text"""
        equations = []
        
        # Pattern for LaTeX equations
        patterns = [
            r'\$\$(.+?)\$\$',  # Display math
            r'\$(.+?)\$',      # Inline math
            r'\\begin{equation}(.+?)\\end{equation}',
            r'\\begin{align}(.+?)\\end{align}',
        ]
        
        for pattern in patterns:
            matches = re.finditer(pattern, text, re.DOTALL)
            for match in matches:
                equations.append({
                    "latex": match.group(1).strip(),
                    "type": "display" if "$$" in match.group(0) else "inline",
                    "position": match.start()
                })
        
        return equations
    
    def convert_latex(self, latex_expr: str) -> str:
        """
        Convert LaTeX to readable format
        
        Args:
            latex_expr: LaTeX expression
            
        Returns:
            Converted string
        """
        # Simple conversion for common patterns
        conversions = {
            r'\\frac{(.+?)}{(.+?)}': r'(\1)/(\2)',
            r'\\sqrt{(.+?)}': r'sqrt(\1)',
            r'\^{(.+?)}': r'^(\1)',
            r'_{(.+?)}': r'_(\1)',
        }
        
        result = latex_expr
        for pattern, replacement in conversions.items():
            result = re.sub(pattern, replacement, result)
        
        return result
    
    def process_tables(self, tables: List[List[str]]) -> List[Dict[str, Any]]:
        """
        Process and structure tables
        
        Args:
            tables: List of table data
            
        Returns:
            Structured table data
        """
        processed = []
        
        for table_idx, table_data in enumerate(tables):
            if not table_data:
                continue
            
            # Detect header row
            header = table_data[0] if table_data else []
            rows = table_data[1:] if len(table_data) > 1 else []
            
            # Convert to dict format
            structured = {
                "table_id": table_idx + 1,
                "headers": header,
                "rows": rows,
                "row_count": len(rows),
                "column_count": len(header),
                "data": [
                    dict(zip(header, row))
                    for row in rows
                ]
            }
            
            processed.append(structured)
        
        return processed


class KnowledgeGraphClient:
    """HTTP client for Knowledge Graph Agent"""
    
    def __init__(self, config: Dict[str, Any]):
        self.config = config
        self.base_url = config["knowledge_graph_api"]["base_url"]
        self.timeout = config["knowledge_graph_api"]["timeout"]
        self.retry_count = config["knowledge_graph_api"]["retry"]
        self.client = httpx.AsyncClient(timeout=self.timeout)
        self.logger = structlog.get_logger()
    
    @retry(
        stop=stop_after_attempt(3),
        wait=wait_exponential(multiplier=1, min=2, max=10)
    )
    async def create_nodes(
        self,
        concepts: List[Dict[str, Any]],
        metadata: Dict[str, Any]
    ) -> List[str]:
        """
        Create concept nodes in knowledge graph
        
        Args:
            concepts: List of concepts
            metadata: File metadata
            
        Returns:
            List of created node IDs
        """
        start_time = datetime.utcnow()
        node_ids = []
        
        try:
            for concept in concepts:
                payload = {
                    "label": "Concept",
                    "properties": {
                        "id": self._generate_concept_id(concept["text"]),
                        "name": concept["text"],
                        "type": concept.get("type", "UNKNOWN"),
                        "confidence": concept.get("confidence", 0.0),
                        "source": metadata.get("file_name", "unknown"),
                        "subject": metadata.get("subject", "general"),
                        "created_at": datetime.utcnow().isoformat()
                    }
                }
                
                response = await self.client.post(
                    f"{self.base_url}/nodes",
                    json=payload
                )
                response.raise_for_status()
                
                result = response.json()
                node_ids.append(result["node_id"])
                
                GRAPH_API_CALLS.labels(operation="create_node", status="success").inc()
            
            duration = (datetime.utcnow() - start_time).total_seconds()
            
            self.logger.info(
                "nodes_created",
                count=len(node_ids),
                duration=duration
            )
            
            return node_ids
            
        except Exception as e:
            GRAPH_API_CALLS.labels(operation="create_node", status="error").inc()
            self.logger.error("node_creation_failed", error=str(e))
            raise
    
    @retry(
        stop=stop_after_attempt(3),
        wait=wait_exponential(multiplier=1, min=2, max=10)
    )
    async def create_relationships(
        self,
        relationships: List[Dict[str, Any]]
    ) -> List[str]:
        """
        Create relationships in knowledge graph
        
        Args:
            relationships: List of relationships
            
        Returns:
            List of created relationship IDs
        """
        start_time = datetime.utcnow()
        rel_ids = []
        
        try:
            for rel in relationships:
                payload = {
                    "from_id": self._generate_concept_id(rel["from_concept"]),
                    "to_id": self._generate_concept_id(rel["to_concept"]),
                    "rel_type": self._map_relationship_type(rel["relationship_type"]),
                    "properties": {
                        "confidence": rel.get("confidence", 0.0),
                        "context": rel.get("context", ""),
                        "created_at": datetime.utcnow().isoformat()
                    }
                }
                
                response = await self.client.post(
                    f"{self.base_url}/relationships",
                    json=payload
                )
                response.raise_for_status()
                
                result = response.json()
                rel_ids.append(result["rel_id"])
                
                GRAPH_API_CALLS.labels(operation="create_relationship", status="success").inc()
            
            duration = (datetime.utcnow() - start_time).total_seconds()
            
            self.logger.info(
                "relationships_created",
                count=len(rel_ids),
                duration=duration
            )
            
            return rel_ids
            
        except Exception as e:
            GRAPH_API_CALLS.labels(operation="create_relationship", status="error").inc()
            self.logger.error("relationship_creation_failed", error=str(e))
            raise
    
    async def close(self):
        """Close HTTP client"""
        await self.client.aclose()
    
    def _generate_concept_id(self, text: str) -> str:
        """Generate unique concept ID"""
        return hashlib.md5(text.lower().encode()).hexdigest()[:16]
    
    def _map_relationship_type(self, rel_type: str) -> str:
        """Map relationship type to graph schema"""
        mapping = {
            "prerequisite_of": "PREREQUISITE_OF",
            "related_to": "BELONGS_TO",
            "example_of": "BELONGS_TO",
            "part_of": "BELONGS_TO",
            "causes": "BELONGS_TO",
            "defines": "BELONGS_TO"
        }
        return mapping.get(rel_type, "BELONGS_TO")


class ContentIngestionAgent:
    """Main Content Ingestion Agent"""
    
    def __init__(self, config_path: str = "config.yaml"):
        self.config = self._load_config(config_path)
        self.model_manager = ModelManager(self.config)
        self.file_processor = FileProcessor(self.config)
        self.graph_client = KnowledgeGraphClient(self.config)
        self.jobs = {}  # Job tracking
        self.logger = structlog.get_logger()
    
    def _load_config(self, config_path: str) -> Dict[str, Any]:
        """Load configuration from YAML"""
        try:
            with open(config_path, 'r') as f:
                config = yaml.safe_load(f)
            logger.info("config_loaded", path=config_path)
            return config
        except Exception as e:
            logger.error("config_load_failed", error=str(e))
            raise
    
    async def initialize(self):
        """Initialize agent"""
        self.logger.info("initializing_agent")
        self.model_manager.load_models()
        self.logger.info("agent_initialized")
    
    async def shutdown(self):
        """Shutdown agent"""
        await self.graph_client.close()
        self.logger.info("agent_shutdown")
    
    async def ingest_file(
        self,
        file_path: Path,
        metadata: Dict[str, Any]
    ) -> Dict[str, Any]:
        """
        Main ingestion pipeline
        
        Args:
            file_path: Path to file
            metadata: File metadata
            
        Returns:
            Ingestion results
        """
        job_id = str(uuid.uuid4())
        file_id = hashlib.md5(str(file_path).encode()).hexdigest()
        
        start_time = datetime.utcnow()
        
        self.jobs[job_id] = {
            "status": ProcessingStatus.PROCESSING,
            "progress": 0.0,
            "message": "Starting ingestion",
            "file_id": file_id,
            "created_at": start_time
        }
        
        try:
            # Step 1: Validate file (10%)
            self.file_processor.validate_file(file_path)
            self._update_job(job_id, 10, "File validated")
            
            # Step 2: Extract content (30%)
            content = self.file_processor.extract_content(file_path)
            self._update_job(job_id, 30, "Content extracted")
            
            # Step 3: Analyze and classify (50%)
            full_text = self._get_full_text(content)
            classification = self.model_manager.classify_content(full_text)
            self._update_job(job_id, 50, "Content classified")
            
            # Step 4: Extract concepts (70%)
            concepts = self.model_manager.extract_concepts(
                full_text,
                method=self.config["extraction"]["concept_extraction_method"]
            )
            self._update_job(job_id, 70, "Concepts extracted")
            
            # Step 5: Detect relationships (80%)
            relationships = []
            if self.config["extraction"]["relationship_detection"]:
                relationships = self.model_manager.detect_relationships(
                    concepts,
                    full_text
                )
            self._update_job(job_id, 80, "Relationships detected")
            
            # Step 6: Create graph nodes (90%)
            metadata.update({
                "file_name": file_path.name,
                "subject": classification["primary_subject"],
                "difficulty": classification["difficulty"],
                "blooms_level": classification["blooms_level"]
            })
            
            node_ids = await self.graph_client.create_nodes(concepts, metadata)
            self._update_job(job_id, 90, "Graph nodes created")
            
            # Step 7: Create relationships (95%)
            rel_ids = []
            if relationships:
                rel_ids = await self.graph_client.create_relationships(relationships)
            self._update_job(job_id, 95, "Relationships created")
            
            # Step 8: Complete (100%)
            duration = (datetime.utcnow() - start_time).total_seconds()
            
            result = {
                "job_id": job_id,
                "file_id": file_id,
                "status": ProcessingStatus.COMPLETED,
                "content": content,
                "classification": classification,
                "concepts": concepts,
                "relationships": relationships,
                "graph_nodes": node_ids,
                "graph_relationships": rel_ids,
                "duration": duration,
                "completed_at": datetime.utcnow().isoformat()
            }
            
            self.jobs[job_id]["status"] = ProcessingStatus.COMPLETED
            self.jobs[job_id]["progress"] = 100.0
            self.jobs[job_id]["message"] = "Ingestion completed"
            self.jobs[job_id]["result"] = result
            
            INGESTION_COUNT.labels(
                status="success",
                file_type=file_path.suffix.lstrip('.')
            ).inc()
            
            self.logger.info(
                "ingestion_completed",
                job_id=job_id,
                duration=duration,
                concepts=len(concepts),
                relationships=len(relationships)
            )
            
            return result
            
        except Exception as e:
            self.jobs[job_id]["status"] = ProcessingStatus.FAILED
            self.jobs[job_id]["message"] = f"Error: {str(e)}"
            
            INGESTION_COUNT.labels(
                status="failed",
                file_type=file_path.suffix.lstrip('.')
            ).inc()
            
            self.logger.error("ingestion_failed", job_id=job_id, error=str(e))
            raise
    
    def _update_job(self, job_id: str, progress: float, message: str):
        """Update job status"""
        if job_id in self.jobs:
            self.jobs[job_id]["progress"] = progress
            self.jobs[job_id]["message"] = message
    
    def _get_full_text(self, content: Dict[str, Any]) -> str:
        """Combine all text from content"""
        texts = []
        for page in content.get("pages", []):
            texts.append(page.get("text", ""))
        return "\n".join(texts)
    
    def get_job_status(self, job_id: str) -> Dict[str, Any]:
        """Get job status"""
        if job_id not in self.jobs:
            raise ValueError(f"Job not found: {job_id}")
        return self.jobs[job_id]
    
    async def get_health_status(self) -> Dict[str, Any]:
        """Get agent health status"""
        return {
            "status": "healthy",
            "timestamp": datetime.utcnow().isoformat(),
            "models_loaded": len(self.model_manager.models) > 0,
            "active_jobs": sum(
                1 for job in self.jobs.values()
                if job["status"] == ProcessingStatus.PROCESSING
            ),
            "version": "1.0.0"
        }


# Global agent instance
agent = None


@asynccontextmanager
async def lifespan(app: FastAPI):
    """FastAPI lifespan handler"""
    global agent
    agent = ContentIngestionAgent()
    await agent.initialize()
    yield
    await agent.shutdown()


# Create FastAPI app
app = FastAPI(
    title="Content Ingestion Agent",
    description="Multi-modal content processing with ML models",
    version="1.0.0",
    lifespan=lifespan
)


@app.post("/ingest", response_model=IngestionResponse, status_code=status.HTTP_202_ACCEPTED)
async def ingest_file_endpoint(
    file: UploadFile = File(...),
    metadata: str = ""
):
    """Upload and process file"""
    try:
        # Parse metadata
        meta_dict = json.loads(metadata) if metadata else {}
        
        # Save uploaded file
        file_path = Path(agent.config["storage"]["path"]) / file.filename
        file_path.parent.mkdir(parents=True, exist_ok=True)
        
        with open(file_path, "wb") as f:
            content = await file.read()
            f.write(content)
        
        # Start ingestion (async)
        job_id = str(uuid.uuid4())
        file_id = hashlib.md5(str(file_path).encode()).hexdigest()
        
        # Queue job
        asyncio.create_task(agent.ingest_file(file_path, meta_dict))
        
        return IngestionResponse(
            job_id=job_id,
            status=ProcessingStatus.PENDING,
            file_id=file_id,
            message="Ingestion started",
            created_at=datetime.utcnow()
        )
        
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=str(e)
        )


@app.get("/status/{job_id}", response_model=StatusResponse)
async def get_status_endpoint(job_id: str):
    """Check processing status"""
    try:
        job = agent.get_job_status(job_id)
        
        return StatusResponse(
            job_id=job_id,
            status=job["status"],
            progress=job["progress"],
            message=job["message"],
            result=job.get("result")
        )
        
    except ValueError as e:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=str(e)
        )


@app.get("/health")
async def health_check():
    """Health check endpoint"""
    health = await agent.get_health_status()
    status_code = status.HTTP_200_OK if health["status"] == "healthy" else status.HTTP_503_SERVICE_UNAVAILABLE
    return JSONResponse(content=health, status_code=status_code)


@app.get("/metrics")
async def metrics():
    """Prometheus metrics"""
    return JSONResponse(
        content=generate_latest().decode("utf-8"),
        media_type="text/plain"
    )


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8001)
